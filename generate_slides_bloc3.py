#!/usr/bin/env python3
"""
Génère le support de soutenance CESIZen — Activité 3 « Déployer et sécuriser »
(PowerPoint 16:9 éditable). Réutilise la charte du support du Bloc 1.
Sortie : ./CESIZen_Soutenance_Bloc3.pptx
Schémas : ./assets_bloc3/*.png
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Charte CESIZen ────────────────────────────────────────────────
GREEN   = RGBColor(6, 198, 86)
GREEN_D = RGBColor(4, 150, 65)
BLUE    = RGBColor(0, 70, 130)
BLUE_L  = RGBColor(230, 240, 250)
YELLOW  = RGBColor(252, 225, 23)
ORANGE  = RGBColor(228, 128, 28)
RED     = RGBColor(200, 16, 46)
GREY_D  = RGBColor(50, 50, 50)
GREY    = RGBColor(110, 110, 110)
WHITE   = RGBColor(255, 255, 255)
INK     = RGBColor(35, 35, 45)

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, 'assets_bloc3')
OUT = os.path.join(HERE, 'CESIZen_Soutenance_Bloc3.pptx')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line is None:
        _no_line(s)
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    return tb, tf


def set_par(p, text, size, color, bold=False, align=PP_ALIGN.LEFT,
            font='Arial', space_after=8):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font


def footer(slide):
    tb, tf = textbox(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35))
    set_par(tf.paragraphs[0],
            "CESIZen · Bloc INFCDAAL3 — Déployer et sécuriser les applications informatiques · Adam Marzuk",
            9, GREY, align=PP_ALIGN.CENTER, space_after=0)


def band_title(slide, text, num=None):
    rect(slide, 0, 0, SW, Inches(1.05), BLUE)
    rect(slide, 0, Inches(1.05), SW, Inches(0.06), GREEN)
    label = f"{num}.  {text}" if num else text
    tb, tf = textbox(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.85),
                     anchor=MSO_ANCHOR.MIDDLE)
    set_par(tf.paragraphs[0], label, 28, WHITE, bold=True, space_after=0)


def bullets(slide, items, x, y, w, h, size=18, gap=10):
    tb, tf = textbox(slide, x, y, w, h)
    for i, it in enumerate(items):
        lvl = 0
        txt = it
        if isinstance(it, tuple):
            txt, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = "•  " if lvl == 0 else "–  "
        set_par(p, prefix + txt, size if lvl == 0 else size - 2,
                INK if lvl == 0 else GREY_D, bold=False, space_after=gap)
        p.level = lvl
        if lvl == 1:
            p.runs[0].font.color.rgb = GREY
    return tb


def image_fit(slide, name, x, y, max_w, max_h, caption=None):
    path = os.path.join(ASSETS, name)
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = max_w / max_h
    if box_ar > ar:
        h = max_h
        w = int(h * ar)
    else:
        w = max_w
        h = int(w / ar)
    left = x + (max_w - w) // 2
    top = y + (max_h - h) // 2
    slide.shapes.add_picture(path, left, top, width=w, height=h)
    if caption:
        tb, tf = textbox(slide, x, y + max_h + Inches(0.02), max_w, Inches(0.35))
        set_par(tf.paragraphs[0], caption, 11, GREY, align=PP_ALIGN.CENTER, space_after=0)


def cards(slide, items, y, h=1.8, top=0.55):
    """items = list of (titre, description, couleur)"""
    n = len(items)
    total_w = 12.0
    gap = 0.25
    cw = (total_w - gap * (n - 1)) / n
    for i, (t, d, col) in enumerate(items):
        x = Inches(0.7 + i * (cw + gap))
        rect(slide, x, Inches(y), Inches(cw), Inches(h), BLUE_L)
        rect(slide, x, Inches(y), Inches(cw), Inches(top), col)
        tb, tf = textbox(slide, x + Inches(0.1), Inches(y + 0.03), Inches(cw - 0.2),
                         Inches(top), MSO_ANCHOR.MIDDLE)
        set_par(tf.paragraphs[0], t, 13, WHITE, bold=True, space_after=0, align=PP_ALIGN.CENTER)
        tb, tf = textbox(slide, x + Inches(0.15), Inches(y + top + 0.1), Inches(cw - 0.3),
                         Inches(h - top - 0.15))
        for j, line in enumerate(d):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            set_par(p, "• " + line, 11.5, INK, space_after=5)


# ════════════════════════════════════════════════════════════════════
# Slide 1 — Titre
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, Inches(2.5), SW, Inches(0.08), GREEN)
tb, tf = textbox(s, Inches(1), Inches(0.9), Inches(11.3), Inches(1.5))
set_par(tf.paragraphs[0], "CESIZen", 60, GREEN, bold=True, align=PP_ALIGN.CENTER, space_after=0)
tb, tf = textbox(s, Inches(1), Inches(2.62), Inches(11.3), Inches(0.9))
set_par(tf.paragraphs[0], "L'application de votre santé mentale", 24, GREY_D,
        align=PP_ALIGN.CENTER, space_after=0)
tb, tf = textbox(s, Inches(1), Inches(3.65), Inches(11.3), Inches(0.9))
set_par(tf.paragraphs[0], "Activité 3 — Déployer et sécuriser les applications informatiques",
        18, BLUE, align=PP_ALIGN.CENTER, space_after=0)
tb, tf = textbox(s, Inches(1), Inches(5.4), Inches(11.3), Inches(1.4))
for i, line in enumerate([
        "Adam Marzuk — CESI École d'Ingénieurs",
        "Bloc INFCDAAL3 — Déployer et sécuriser les applications informatiques",
        "Commanditaire : Ministère des Solidarités et de la Santé"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    set_par(p, line, 15, GREY, align=PP_ALIGN.CENTER, space_after=4)

# ════════════════════════════════════════════════════════════════════
# Slide 2 — Sommaire
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Sommaire")
bullets(s, [
    "Contexte et périmètre de l'Activité 3",
    "Plan de déploiement : architecture et environnements",
    "Versioning et intégration/livraison continues (Gitea · Coolify)",
    "Plan de maintenance : ticketing, engagements de service, veille",
    "Plan de sécurisation : vulnérabilités, risques, actions",
    "Données personnelles, RGPD et gestion de crise",
    "Démonstration des outils et conclusion",
], Inches(1.2), Inches(1.5), Inches(11), Inches(5.2), size=22, gap=13)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 3 — Contexte & périmètre
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Contexte et périmètre de l'Activité 3", 1)
bullets(s, [
    "CESIZen : plateforme grand public de santé mentale (Ministère).",
    "Après la conception (Act. 1) et le développement/test (Act. 2), l'Activité 3 traite du déploiement et de la sécurisation.",
    "Prototype déployé : API Laravel + web Next.js + mobile Expo + PostgreSQL.",
], Inches(0.8), Inches(1.4), Inches(11.7), Inches(2.0), size=18)
tb, tf = textbox(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.5))
set_par(tf.paragraphs[0], "3 livrables", 16, GREEN_D, bold=True, space_after=0)
cards(s, [
    ("Plan de déploiement", ["Environnements", "Versioning", "Automatisations"], GREEN),
    ("Plan de maintenance", ["Ticketing Gitea", "SLA & méthodologie", "Veille techno"], BLUE),
    ("Plan de sécurisation", ["Risques & criticité", "Actions correctives", "RGPD & crise"], GREEN_D),
], y=4.3, h=2.0)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 4 — Architecture technique
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Architecture technique de la solution", 2)
bullets(s, [
    "Architecture découplée : une API REST centrale, deux clients (web + mobile).",
    ("Back-end : Laravel 12 (PHP 8.4) — API REST JSON, authentification JWT.", 1),
    ("Base de données : PostgreSQL 16 (ACID) ; Redis en développement.", 1),
    ("Front-end web : Next.js 16 / React 19, design system de l'État (DSFR).", 1),
    ("Mobile : React Native / Expo 54 (iOS + Android, Mobile First).", 1),
    "Conteneurisation Docker : environnements reproductibles et isolés.",
    "Une seule logique métier → non-duplication, testabilité, évolutivité.",
], Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.0), size=17, gap=9)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 5 — Plan de déploiement : 3 environnements + chaîne
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Plan de déploiement — environnements et chaîne CI/CD", 2)
image_fit(s, "archi_deploiement.png", Inches(0.5), Inches(1.25), Inches(12.3), Inches(5.5),
          caption="Développement (local) → Test/Intégration (Gitea Actions) → Production (Coolify, HTTPS)")
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 6 — Versioning & CI/CD
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Versioning et intégration/livraison continues", 3)
bullets(s, [
    "Gitea : dépôt Git, ticketing et CI/CD réunis dans un outil souverain auto-hébergé.",
    "Stratégie de branches : main (prod, protégée) · develop · feature/fix/hotfix.",
    "Conventions : Conventional Commits, revue par pull request, tags de version.",
], Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.9), size=17)
tb, tf = textbox(s, Inches(0.8), Inches(3.55), Inches(11.7), Inches(0.5))
set_par(tf.paragraphs[0], "Deux automatisations complémentaires", 16, GREEN_D, bold=True, space_after=0)
cards(s, [
    ("Intégration continue — Gitea Actions", [
        "29 tests à chaque push",
        "Lint (Pint, ESLint)",
        "Merge bloqué si CI rouge"], ORANGE),
    ("Déploiement continu — Coolify", [
        "Webhook sur merge main",
        "Build images + migrations",
        "HTTPS, health checks, restart"], BLUE),
], y=4.15, h=2.1)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 7 — Maintenance : ticketing
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Plan de maintenance — ticketing et cycle de vie", 4)
image_fit(s, "cycle_ticket.png", Inches(0.5), Inches(1.25), Inches(12.3), Inches(5.5),
          caption="Gitea Issues : gabarits incident/évolution, labels, jalons, Kanban · SLA du contrat")
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 8 — Maintenance : méthodo + veille
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Maintenance corrective, évolutive et veille", 4)
tb, tf = textbox(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5))
set_par(tf.paragraphs[0], "Méthodologie", 16, GREEN_D, bold=True, space_after=0)
bullets(s, [
    "Corrective : ticket → qualification (gravité) → SLA → correction → revue (PR + CI) → clôture validée.",
    "Évolutive : analyse (documentation, délai, coût) → validation client → planification par version.",
    "Pilotage : tableaux de bord Gitea (Kanban, jalons, indicateurs de délais).",
], Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.3), size=16, gap=9)
tb, tf = textbox(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.5))
set_par(tf.paragraphs[0], "Veille technologique (pérennité)", 16, GREEN_D, bold=True, space_after=0)
bullets(s, [
    "Sécurité : CVE, CERT-FR, OWASP, alertes de dépendances (composer/npm audit en CI).",
    "Frameworks & dépendances : notes de version, robot de mises à jour testées par la CI.",
    "Réglementation : CNIL, ANSSI (RGPD, HDS).",
], Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.9), size=16, gap=9)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 9 — Sécurisation : vulnérabilités OWASP
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Plan de sécurisation — vulnérabilités (OWASP)", 5)
bullets(s, [
    "Contrôle d'accès : RBAC par middleware, garde-fous admin.",
    "Cryptographie : HTTPS/TLS, bcrypt (coût 12), mots de passe masqués et hors logs.",
    "Injection : ORM paramétré (Eloquent) + validation par Form Requests.",
    "Authentification : JWT signé, TTL 60 min, blacklist, rate limiting différencié.",
    "Journalisation : piste d'audit immuable (IP, anciennes/nouvelles valeurs).",
    "Configuration : APP_DEBUG=false, secrets hors dépôt, CORS en liste blanche.",
], Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.0), size=17, gap=11)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 10 — Matrice de risques
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Sécurisation — matrice de criticité des risques", 5)
image_fit(s, "matrice_risques.png", Inches(2.4), Inches(1.25), Inches(8.5), Inches(5.5),
          caption="Criticité = gravité × probabilité · priorisation des actions (R1, R7, R2 prioritaires)")
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 11 — Actions préventives / correctives
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Sécurisation — actions préventives et correctives", 5)
cards(s, [
    ("Prévention", [
        "Validation & ORM (injection)",
        "MFA admin, mots de passe robustes",
        "audit dépendances en CI",
        "Sauvegardes automatisées"], GREEN_D),
    ("Chiffrement", [
        "HTTPS/TLS en transit",
        "bcrypt + AES (APP_KEY)",
        "Notes de santé chiffrées (prévu)",
        "JWT signés"], BLUE),
    ("Correction", [
        "Correctif prioritaire + test",
        "Révocation / blacklist token",
        "Redéploiement automatisé",
        "Plan de reprise (PRA)"], ORANGE),
], y=1.6, h=4.6, top=0.6)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 12 — RGPD
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Données personnelles et RGPD", 6)
bullets(s, [
    "Données de santé mentale = données sensibles (RGPD art. 9).",
    "Consentement explicite obligatoire à l'inscription, tracé en base.",
    "Droits implémentés : accès et portabilité (export JSON), effacement (anonymisation + soft delete).",
    "Minimisation des données ; hébergement UE, aucun transfert hors UE ; cible HDS en production.",
], Inches(0.8), Inches(1.45), Inches(11.7), Inches(2.4), size=17, gap=10)
rect(s, Inches(0.8), Inches(4.25), Inches(11.7), Inches(2.3), BLUE_L)
rect(s, Inches(0.8), Inches(4.25), Inches(3.6), Inches(0.55), GREEN_D)
tb, tf = textbox(s, Inches(0.95), Inches(4.28), Inches(3.4), Inches(0.5), MSO_ANCHOR.MIDDLE)
set_par(tf.paragraphs[0], "Traçabilité (art. 32)", 15, WHITE, bold=True, space_after=0)
bullets(s, [
    "Piste d'audit immuable : action, IP, anciennes/nouvelles valeurs — le mot de passe n'est jamais journalisé.",
    "Chiffrement en transit, hachage bcrypt, contrôle d'accès par rôle.",
    "Export et anonymisation testés automatiquement (tests RGPD).",
], Inches(1.0), Inches(4.95), Inches(11.3), Inches(1.6), size=14, gap=6)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 13 — Gestion de crise & escalade
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Gestion de crise et escalade d'information", 6)
levels = [
    ("N1 — Détection", "Astreinte technique", "Qualifier, contenir, ouvrir un ticket sécurité", GREEN),
    ("N2 — Confirmation", "Responsable technique", "Isoler, préserver les preuves, évaluer l'impact", YELLOW),
    ("N3 — Crise", "Chef de projet + DPO", "Notification CNIL < 72 h, information des personnes", ORANGE),
    ("N4 — Direction", "Direction + Ministère", "Communication officielle, coordination ANSSI", RED),
]
y = 1.55
for titre, resp, act, col in levels:
    rect(s, Inches(0.8), Inches(y), Inches(11.7), Inches(1.15), BLUE_L)
    rect(s, Inches(0.8), Inches(y), Inches(0.16), Inches(1.15), col)
    tb, tf = textbox(s, Inches(1.1), Inches(y + 0.05), Inches(3.0), Inches(1.05), MSO_ANCHOR.MIDDLE)
    set_par(tf.paragraphs[0], titre, 15, col if col != YELLOW else GREY_D, bold=True, space_after=2)
    p = tf.add_paragraph(); set_par(p, resp, 12, GREY_D, space_after=0)
    tb, tf = textbox(s, Inches(4.3), Inches(y + 0.05), Inches(8.0), Inches(1.05), MSO_ANCHOR.MIDDLE)
    set_par(tf.paragraphs[0], act, 14, INK, space_after=0)
    y += 1.32
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 14 — Démo & conclusion
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Démonstration et conclusion", 7)
tb, tf = textbox(s, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.5))
set_par(tf.paragraphs[0], "Démonstration des outils (en direct)", 17, GREEN_D, bold=True, space_after=0)
bullets(s, [
    "Versioning : arbre des commits, branches et tags sur Gitea.",
    "Intégration continue : pull request → 29 tests verts (Gitea Actions).",
    "Ticketing : création d'une issue (gabarit incident), labels, jalon, Kanban.",
    "Déploiement : merge → Coolify → site en ligne (cesizen.cleanows.fr).",
], Inches(0.8), Inches(1.85), Inches(11.7), Inches(2.4), size=16, gap=9)
tb, tf = textbox(s, Inches(0.8), Inches(4.35), Inches(11.7), Inches(0.5))
set_par(tf.paragraphs[0], "En synthèse", 17, GREEN_D, bold=True, space_after=0)
bullets(s, [
    "Chaîne de déploiement complète, outillée et réellement en production.",
    "Maintenance cadrée (ticketing, SLA, veille) et sécurité analysée par les risques + RGPD.",
], Inches(0.8), Inches(4.85), Inches(11.7), Inches(1.4), size=16, gap=8)
rect(s, 0, Inches(6.55), SW, Inches(0.95), BLUE)
tb, tf = textbox(s, Inches(0.8), Inches(6.57), Inches(11.7), Inches(0.9), MSO_ANCHOR.MIDDLE)
set_par(tf.paragraphs[0], "Merci de votre attention — Questions / Réponses", 20, WHITE,
        bold=True, align=PP_ALIGN.CENTER, space_after=0)

# ─── Numéros de page (toutes les slides sauf la page de titre) ─────
for idx, slide in enumerate(prs.slides):
    if idx == 0:
        continue  # pas de numéro sur la slide de titre
    tb, tf = textbox(slide, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.35),
                     anchor=MSO_ANCHOR.MIDDLE)
    set_par(tf.paragraphs[0], str(idx + 1), 10, GREY, bold=True,
            align=PP_ALIGN.RIGHT, space_after=0)

prs.save(OUT)
print(f"Support de soutenance généré : {OUT}")
print(f"Nombre de slides : {len(prs.slides._sldIdLst)}")
