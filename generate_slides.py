#!/usr/bin/env python3
"""
Génère le support de soutenance CESIZen (PowerPoint 16:9 éditable).

Réutilise les schémas PNG produits par generate_diagrams.py.
Sortie : ./CESIZen_Soutenance.pptx
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Charte CESIZen ────────────────────────────────────────────────
GREEN   = RGBColor(6, 198, 86)
GREEN_D = RGBColor(4, 150, 65)
BLUE    = RGBColor(0, 70, 130)
BLUE_L  = RGBColor(230, 240, 250)
YELLOW  = RGBColor(252, 225, 23)
GREY_D  = RGBColor(50, 50, 50)
GREY    = RGBColor(110, 110, 110)
WHITE   = RGBColor(255, 255, 255)
INK     = RGBColor(35, 35, 45)

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, 'assets_dossier')
OUT = os.path.join(HERE, 'CESIZen_Soutenance.pptx')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line is None:
        _no_line(s)
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    return tb, tf


def set_par(p, text, size, color, bold=False, align=PP_ALIGN.LEFT,
            font='Arial', space_after=8):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font


def footer(slide):
    tb, tf = textbox(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35))
    set_par(tf.paragraphs[0],
            "CESIZen · Bloc INFCDAAL1 — Concevoir les solutions logicielles · Adam Marzuk",
            9, GREY, align=PP_ALIGN.CENTER, space_after=0)


def band_title(slide, text, num=None):
    rect(slide, 0, 0, SW, Inches(1.05), BLUE)
    rect(slide, 0, Inches(1.05), SW, Inches(0.06), GREEN)
    label = f"{num}.  {text}" if num else text
    tb, tf = textbox(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.85),
                     anchor=MSO_ANCHOR.MIDDLE)
    set_par(tf.paragraphs[0], label, 30, WHITE, bold=True, space_after=0)


def bullets(slide, items, x, y, w, h, size=18, gap=10):
    tb, tf = textbox(slide, x, y, w, h)
    for i, it in enumerate(items):
        lvl = 0
        txt = it
        if isinstance(it, tuple):
            txt, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = "•  " if lvl == 0 else "–  "
        set_par(p, prefix + txt, size if lvl == 0 else size - 2,
                INK if lvl == 0 else GREY_D, bold=False, space_after=gap)
        p.level = lvl
        if lvl == 1:
            p.runs[0].font.color.rgb = GREY
    return tb


def image_fit(slide, name, x, y, max_w, max_h, caption=None):
    path = os.path.join(ASSETS, name)
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = max_w / max_h
    if box_ar > ar:
        h = max_h
        w = int(h * ar)
    else:
        w = max_w
        h = int(w / ar)
    left = x + (max_w - w) // 2
    top = y + (max_h - h) // 2
    slide.shapes.add_picture(path, left, top, width=w, height=h)
    if caption:
        tb, tf = textbox(slide, x, y + max_h + Inches(0.02), max_w, Inches(0.35))
        set_par(tf.paragraphs[0], caption, 11, GREY, align=PP_ALIGN.CENTER, space_after=0)


# ════════════════════════════════════════════════════════════════════
# Slide 1 — Titre
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, Inches(2.5), SW, Inches(0.08), GREEN)
tb, tf = textbox(s, Inches(1), Inches(0.9), Inches(11.3), Inches(1.5))
set_par(tf.paragraphs[0], "CESIZen", 60, GREEN, bold=True, align=PP_ALIGN.CENTER, space_after=0)
tb, tf = textbox(s, Inches(1), Inches(2.65), Inches(11.3), Inches(0.9))
set_par(tf.paragraphs[0], "L'application de votre santé mentale", 26, GREY_D,
        align=PP_ALIGN.CENTER, space_after=0)
tb, tf = textbox(s, Inches(1), Inches(3.7), Inches(11.3), Inches(0.7))
set_par(tf.paragraphs[0], "Cahier des charges & conception logicielle", 18, BLUE,
        align=PP_ALIGN.CENTER, space_after=0)
tb, tf = textbox(s, Inches(1), Inches(5.4), Inches(11.3), Inches(1.4))
for i, line in enumerate([
        "Adam Marzuk — CESI École d'Ingénieurs",
        "Bloc INFCDAAL1 — Concevoir les solutions logicielles",
        "Commanditaire : Ministère des Solidarités et de la Santé"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    set_par(p, line, 15, GREY, align=PP_ALIGN.CENTER, space_after=4)

# ════════════════════════════════════════════════════════════════════
# Slide 2 — Sommaire
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Sommaire")
bullets(s, [
    "Contexte, objectifs et enjeux",
    "Parties prenantes et acteurs",
    "Besoins reformulés et priorisation",
    "Réponse fonctionnelle (Comptes · Informations · Tracker)",
    "Modélisation : cas d'usage UML, MCD, architecture MVC",
    "Données personnelles et sensibles (RGPD)",
    "Prototype et perspectives",
], Inches(1.2), Inches(1.5), Inches(11), Inches(5.2), size=22, gap=14)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 3 — Contexte & enjeux
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Contexte, objectifs et enjeux", 1)
bullets(s, [
    "Commanditaire : Ministère des Solidarités et de la Santé.",
    "Enjeu : une plateforme grand public pour prévenir le stress et accompagner la santé mentale.",
    "Constat : 22 % des Français touchés par un trouble psychique au cours de leur vie (Santé Publique France, 2023).",
    "Frein : stigmatisation et méconnaissance limitent l'accès au soin.",
], Inches(0.8), Inches(1.4), Inches(11.7), Inches(2.6), size=18)
# encadré objectifs
rect(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.4), BLUE_L)
rect(s, Inches(0.8), Inches(4.2), Inches(3.4), Inches(0.55), BLUE)
tb, tf = textbox(s, Inches(0.95), Inches(4.25), Inches(3.2), Inches(0.5), MSO_ANCHOR.MIDDLE)
set_par(tf.paragraphs[0], "4 objectifs stratégiques", 15, WHITE, bold=True, space_after=0)
bullets(s, [
    "Informer : contenu éducatif fiable sur la santé mentale.",
    "Suivre : permettre à chacun de tracer son humeur dans le temps.",
    "Orienter : diriger les personnes en détresse (numéros d'urgence).",
    "Protéger : confidentialité absolue des données émotionnelles (RGPD + HDS).",
], Inches(1.0), Inches(4.9), Inches(11.3), Inches(1.6), size=15, gap=4)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 4 — Parties prenantes & acteurs
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Parties prenantes et acteurs", 2)
tb, tf = textbox(s, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.5))
set_par(tf.paragraphs[0], "Parties prenantes du projet", 16, GREEN_D, bold=True, space_after=0)
bullets(s, [
    "Commanditaire : Ministère (pilotage, exigences santé publique et RGPD).",
    "Bénéficiaires : grand public, personnes en quête de mieux-être psychique.",
    "Exploitants : administrateurs et professionnels de santé (modération, contenus).",
], Inches(0.8), Inches(1.85), Inches(11.7), Inches(2.0), size=17)
tb, tf = textbox(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.5))
set_par(tf.paragraphs[0], "3 acteurs du système", 16, GREEN_D, bold=True, space_after=0)
cards = [("Visiteur anonyme", "Consulte les informations, crée un compte.", GREEN),
         ("Utilisateur connecté", "Journal d'émotions, rapports, profil, RGPD.", BLUE),
         ("Administrateur", "Gère utilisateurs, contenus et émotions.", GREEN_D)]
for i, (t, d, col) in enumerate(cards):
    x = Inches(0.8 + i * 4.0)
    rect(s, x, Inches(4.55), Inches(3.7), Inches(1.8), BLUE_L)
    rect(s, x, Inches(4.55), Inches(3.7), Inches(0.55), col)
    tb, tf = textbox(s, x + Inches(0.1), Inches(4.58), Inches(3.5), Inches(0.5), MSO_ANCHOR.MIDDLE)
    set_par(tf.paragraphs[0], t, 14, WHITE, bold=True, space_after=0)
    tb, tf = textbox(s, x + Inches(0.15), Inches(5.2), Inches(3.4), Inches(1.0))
    set_par(tf.paragraphs[0], d, 13, INK, space_after=0)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 5 — Besoins & périmètre
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Besoins reformulés et périmètre", 3)
cols = [
    ("2 modules OBLIGATOIRES", GREEN, [
        "Comptes utilisateurs : inscription JWT, profil, reset, admin, RGPD.",
        "Informations : pages de contenu santé mentale, CRUD admin.",
    ]),
    ("1 module AU CHOIX (retenu)", BLUE, [
        "Tracker d'émotions : journal, saisies, rapports périodiques.",
        "Hiérarchie : 7 émotions de base + 28 sous-émotions.",
    ]),
    ("Modules ÉCARTÉS", GREY, [
        "Diagnostic de stress.",
        "Cohérence cardiaque.",
        "Activités de détente.",
    ]),
]
for i, (t, col, items) in enumerate(cols):
    x = Inches(0.7 + i * 4.05)
    rect(s, x, Inches(1.5), Inches(3.85), Inches(0.6), col)
    tb, tf = textbox(s, x + Inches(0.1), Inches(1.53), Inches(3.65), Inches(0.55), MSO_ANCHOR.MIDDLE)
    set_par(tf.paragraphs[0], t, 13, WHITE, bold=True, space_after=0)
    bullets(s, items, x + Inches(0.05), Inches(2.3), Inches(3.85), Inches(3.5), size=13, gap=8)
tb, tf = textbox(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.8))
set_par(tf.paragraphs[0],
        "Périmètre maîtrisé : 3 modules livrés en qualité plutôt que 6 modules superficiels.",
        15, GREEN_D, bold=True, align=PP_ALIGN.CENTER, space_after=0)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 6 — Priorisation
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Critères de priorisation et pondération", 3)
bullets(s, [
    "5 critères pondérés, total sur 15 : Complexité, Valeur Ministère, Valeur Utilisateur, Nécessité, Interdépendance.",
], Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.8), size=16)
# mini-tableau top priorités
rows = [
    ("Fonctionnalité", "Score /15", True),
    ("Journal d'émotions · Rapports · Diagnostic de stress", "10", False),
    ("Cryptage · Inscription · Cohérence cardiaque · Ajout tracker", "9", False),
    ("Conformité RGPD", "8", False),
    ("Réinitialisation mot de passe", "7", False),
]
ty = 2.3
for i, (a, b, head) in enumerate(rows):
    col = BLUE if head else (BLUE_L if i % 2 else WHITE)
    rect(s, Inches(0.8), Inches(ty + i * 0.62), Inches(9.4), Inches(0.62), col,
         line=RGBColor(210, 220, 230))
    rect(s, Inches(10.2), Inches(ty + i * 0.62), Inches(2.3), Inches(0.62), col,
         line=RGBColor(210, 220, 230))
    tb, tf = textbox(s, Inches(0.95), Inches(ty + i * 0.62), Inches(9.1), Inches(0.62), MSO_ANCHOR.MIDDLE)
    set_par(tf.paragraphs[0], a, 13, WHITE if head else INK, bold=head, space_after=0)
    tb, tf = textbox(s, Inches(10.2), Inches(ty + i * 0.62), Inches(2.3), Inches(0.62), MSO_ANCHOR.MIDDLE)
    set_par(tf.paragraphs[0], b, 13, WHITE if head else GREEN_D, bold=True,
            align=PP_ALIGN.CENTER, space_after=0)
tb, tf = textbox(s, Inches(0.8), Inches(6.05), Inches(11.9), Inches(0.8))
set_par(tf.paragraphs[0],
        "→ Le Tracker cumule 4 fonctionnalités notées 9-10/15 : module au choix le plus rentable.",
        15, GREEN_D, bold=True, space_after=0)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 7 — Réponse fonctionnelle Comptes
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Réponse fonctionnelle — Comptes utilisateurs", 4)
bullets(s, [
    "Inscription et connexion sécurisées par JWT (tymon/jwt-auth).",
    "Consentement RGPD obligatoire, non pré-coché, à l'inscription.",
    "Gestion du profil : édition, changement de mot de passe.",
    "Espace admin : CRUD utilisateurs, activation/désactivation, rôles.",
    "RGPD : export des données (JSON), anonymisation, soft delete + purge 30 j.",
    "Mot de passe : bcrypt 12 rounds, min. 8 caractères.",
], Inches(0.8), Inches(1.5), Inches(8.0), Inches(5.2), size=17, gap=12)
image_fit(s, "mockup_inscription.png", Inches(9.1), Inches(1.4), Inches(3.7), Inches(5.2),
          caption="Écran d'inscription (mobile)")
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 8 — Réponse fonctionnelle Informations
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Réponse fonctionnelle — Informations", 4)
bullets(s, [
    "Pages de contenu santé mentale (feeds) accessibles au grand public.",
    "Page liste (cartes : titre, extrait, image, date) et page détail (HTML rendu).",
    "Back-office admin : création / édition / publication / suppression.",
    "Slug unique généré automatiquement à partir du titre.",
    "Champ booléen est_publie : brouillon (invisible) ou publié (visible).",
    "Modération obligatoire : un administrateur valide avant publication.",
    "Tri par ordre numérique pour mettre en avant un contenu.",
], Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.0), size=18, gap=12)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 9 — Réponse fonctionnelle Tracker
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Réponse fonctionnelle — Tracker d'émotions", 4)
bullets(s, [
    "Wizard d'ajout en 3 étapes : émotion de base → sous-émotion + intensité → date + note.",
    "Journal de bord : consultation, modification, suppression des saisies.",
    "Rapports analytiques sur 4 périodes : semaine, mois, trimestre, année.",
    "Statistiques : intensité moyenne, émotion dominante, répartition, évolution.",
    "Hiérarchie configurable par l'admin : 7 émotions + 28 sous-émotions.",
    "Web (graphiques Recharts) et mobile (Expo) synchronisés via l'API.",
], Inches(0.7), Inches(1.45), Inches(7.0), Inches(5.3), size=16, gap=10)
image_fit(s, "mockup_wizard.png", Inches(7.9), Inches(1.35), Inches(2.5), Inches(5.3))
image_fit(s, "mockup_rapports.png", Inches(10.5), Inches(2.4), Inches(2.7), Inches(3.2))
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 10 — Cas d'usage UML
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Modélisation — Cas d'usage UML", 5)
image_fit(s, "usecase.png", Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.4),
          caption="3 acteurs · cas d'usage groupés par module · modules écartés non représentés")
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 11 — MCD
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Modélisation — MCD (Merise)", 5)
image_fit(s, "mcd.png", Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.4),
          caption="8 entités · cardinalités Merise · hiérarchie d'émotions (auto-référence)")
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 12 — Architecture MVC
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Modélisation — Architecture MVC", 5)
image_fit(s, "mvc.png", Inches(0.6), Inches(1.25), Inches(8.6), Inches(5.5))
bullets(s, [
    "API REST Laravel 11 = Contrôleur central.",
    "Eloquent ORM + PostgreSQL 16 = Modèle.",
    "Next.js (web) + React Native (mobile) = Vues.",
    "Une seule logique métier, deux clients.",
    "Sécurité centralisée : FormRequest, middleware, JWT.",
    "Stack open-source, conteneurisée (Docker).",
], Inches(9.3), Inches(1.5), Inches(3.5), Inches(5.0), size=14, gap=12)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 13 — RGPD & données sensibles
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Données personnelles et sensibles (RGPD)", 6)
bullets(s, [
    "Données de santé mentale = données sensibles (RGPD art. 9).",
    "Base légale : consentement explicite ; finalités documentées.",
    "Droits CNIL implémentés : accès, rectification, effacement, portabilité (export JSON).",
    "Anonymisation et soft delete + purge à 30 jours.",
], Inches(0.8), Inches(1.45), Inches(11.7), Inches(2.4), size=17, gap=10)
rect(s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.5), BLUE_L)
rect(s, Inches(0.8), Inches(4.1), Inches(3.6), Inches(0.55), GREEN_D)
tb, tf = textbox(s, Inches(0.95), Inches(4.13), Inches(3.4), Inches(0.5), MSO_ANCHOR.MIDDLE)
set_par(tf.paragraphs[0], "Mesures de sécurité", 15, WHITE, bold=True, space_after=0)
bullets(s, [
    "Chiffrement en transit (HTTPS/TLS 1.3) et au repos (AES-256).",
    "Authentification JWT signée, expiration 60 min, throttle login.",
    "Hébergement certifié HDS (Scaleway / OVH Healthcare).",
    "Journalisation : table audits trace toute action admin sensible.",
], Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.7), size=14, gap=6)
footer(s)

# ════════════════════════════════════════════════════════════════════
# Slide 14 — Prototype & perspectives
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
band_title(s, "Prototype et perspectives", 7)
tb, tf = textbox(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.5))
set_par(tf.paragraphs[0], "Un prototype fonctionnel accompagne le dossier", 17, GREEN_D,
        bold=True, space_after=0)
bullets(s, [
    "Backend Laravel 11 (API REST JWT) + PostgreSQL 16.",
    "Frontend web Next.js 16 et application mobile Expo (iOS / Android).",
    "Conteneurisation Docker Compose, dépôt Git.",
], Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.9), size=16)
tb, tf = textbox(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.5))
set_par(tf.paragraphs[0], "Perspectives d'évolution", 17, GREEN_D, bold=True, space_after=0)
bullets(s, [
    "Activer les modules écartés (Diagnostic, Cohérence cardiaque, Activités) sans réécrire le modèle.",
    "Open Data : base anonymisée pour études épidémiologiques.",
    "Interopérabilité : Mon Espace Santé (FHIR) ; internationalisation (i18n).",
], Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.7), size=16)
rect(s, 0, Inches(6.6), SW, Inches(0.9), BLUE)
tb, tf = textbox(s, Inches(0.8), Inches(6.62), Inches(11.7), Inches(0.85), MSO_ANCHOR.MIDDLE)
set_par(tf.paragraphs[0], "Merci de votre attention — Questions / Réponses", 20, WHITE,
        bold=True, align=PP_ALIGN.CENTER, space_after=0)

# ════════════════════════════════════════════════════════════════════
# Annexes (backup Q/R) — séquences
# ════════════════════════════════════════════════════════════════════
for name, title, cap in [
    ("seq_login.png", "Annexe — Séquence d'authentification JWT",
     "Validation · vérification bcrypt · émission du token JWT"),
    ("seq_saisie.png", "Annexe — Séquence d'ajout d'une saisie",
     "Wizard 3 étapes · validation · persistance"),
]:
    s = prs.slides.add_slide(BLANK)
    band_title(s, title)
    image_fit(s, name, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.4), caption=cap)
    footer(s)

prs.save(OUT)
print(f"Support de soutenance généré : {OUT}")
print(f"Nombre de slides : {len(prs.slides._sldIdLst)}")
